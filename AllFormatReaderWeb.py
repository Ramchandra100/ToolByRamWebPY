import streamlit as st
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import io
import re
from datetime import datetime
import base64

def run():
    """Main function for Universal File Reader"""
    
    # File type icons and descriptions
    file_types = {
        'pdf': {'icon': '📄', 'name': 'PDF Document'},
        'docx': {'icon': '📝', 'name': 'Word Document'},
        'pptx': {'icon': '📊', 'name': 'PowerPoint'},
        'xlsx': {'icon': '📈', 'name': 'Excel Spreadsheet'},
        'csv': {'icon': '📊', 'name': 'CSV File'},
        'txt': {'icon': '📃', 'name': 'Text File'},
        'jpg': {'icon': '🖼️', 'name': 'JPEG Image'},
        'png': {'icon': '🖼️', 'name': 'PNG Image'}
    }
    
    # Display supported formats in a nice grid
    st.markdown("### 📋 Supported Formats")
    cols = st.columns(4)
    for i, (ext, info) in enumerate(file_types.items()):
        with cols[i % 4]:
            st.markdown(f"{info['icon']} **.{ext}**  \n{info['name']}")
    
    st.divider()
    
    # File upload section
    st.markdown("### 📤 Upload Files")
    
    # Create a container for file upload
    upload_container = st.container()
    
    with upload_container:
        uploaded_files = st.file_uploader(
            "Choose files to view",
            type=list(file_types.keys()),
            accept_multiple_files=True,
            help="Select one or more files to view their contents"
        )
    
    # Process uploaded files
    if uploaded_files:
        st.markdown("### 📖 File Viewer")
        
        # Create tabs for each file
        file_names = [f.name for f in uploaded_files]
        tabs = st.tabs([f"📄 {name[:20]}..." if len(name) > 20 else f"📄 {name}" for name in file_names])
        
        for idx, (tab, uploaded_file) in enumerate(zip(tabs, uploaded_files)):
            with tab:
                file_type = uploaded_file.name.split('.')[-1].lower()
                
                # File info header
                col1, col2, col3 = st.columns([2, 1, 1])
                with col1:
                    st.markdown(f"**File:** {uploaded_file.name}")
                with col2:
                    st.markdown(f"**Size:** {uploaded_file.size / 1024:.1f} KB")
                with col3:
                    st.markdown(f"**Type:** {file_types.get(file_type, {}).get('icon', '📁')} {file_type.upper()}")
                
                st.divider()
                
                # Reset file pointer to beginning
                uploaded_file.seek(0)
                
                # Create a container for the file content with border
                content_container = st.container()
                
                with content_container:
                    try:
                        # PDF Reader
                        if file_type == 'pdf':
                            display_pdf(uploaded_file, idx)
                        
                        # Word Document
                        elif file_type == 'docx':
                            display_docx(uploaded_file, idx)
                        
                        # PowerPoint
                        elif file_type == 'pptx':
                            display_pptx(uploaded_file, idx)
                        
                        # Excel
                        elif file_type == 'xlsx':
                            display_excel(uploaded_file, idx)
                        
                        # CSV
                        elif file_type == 'csv':
                            display_csv(uploaded_file, idx)
                        
                        # Text File
                        elif file_type == 'txt':
                            display_text(uploaded_file, idx)
                        
                        # Images
                        elif file_type in ['jpg', 'jpeg', 'png']:
                            display_image(uploaded_file, idx)
                        
                    except Exception as e:
                        st.error(f"❌ Error reading file: {str(e)}")
                        st.exception(e)
    
    else:
        # Show example when no files are uploaded
        st.info("👆 Upload files above to view their contents")
        
        # Demo section
        with st.expander("📋 Quick Guide", expanded=True):
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("""
                **✨ What you can do:**
                - 📄 View PDF documents with page navigation
                - 📝 Read Word documents with formatting
                - 📊 Browse Excel/CSV files
                - 🖼️ Preview images and view metadata
                - 📃 Read text files with line numbers
                """)
            
            with col2:
                st.markdown("""
                **🎯 Pro Tips:**
                - Upload multiple files at once
                - Use tabs to switch between files
                - Search within documents
                - Download extracted content
                - View file metadata
                """)

def display_pdf(file, tab_id):
    """Display PDF file with navigation"""
    try:
        reader = PdfReader(file)
        total_pages = len(reader.pages)
        
        # PDF metadata
        if reader.metadata:
            with st.expander("📋 Document Properties"):
                for key, value in reader.metadata.items():
                    if value:
                        st.text(f"{key.replace('/', '')}: {value}")
        
        # Page navigation
        col1, col2, col3 = st.columns([2, 1, 1])
        with col1:
            page_num = st.number_input(
                "📄 Page",
                min_value=1,
                max_value=total_pages,
                value=1,
                key=f"pdf_page_{tab_id}"
            )
        with col2:
            st.metric("Total Pages", total_pages)
        with col3:
            search_term = st.text_input("🔍 Search", "", key=f"pdf_search_{tab_id}", placeholder="Search text...")
        
        # Extract and display text
        page = reader.pages[page_num - 1]
        text = page.extract_text()
        
        # Highlight search results
        if search_term:
            text = re.sub(
                f'({re.escape(search_term)})',
                r'**\1**',
                text,
                flags=re.IGNORECASE
            )
        
        # Display text
        st.text_area(
            f"Page {page_num} Content",
            text if text else "(No text found on this page)",
            height=400,
            key=f"pdf_content_{tab_id}"
        )
        
        # Download button
        if st.button(f"📥 Download Page {page_num}", key=f"pdf_dl_{tab_id}"):
            st.download_button(
                "Click to Download",
                text,
                f"page_{page_num}.txt",
                "text/plain"
            )
    
    except Exception as e:
        st.error(f"Error reading PDF: {str(e)}")

def display_docx(file, tab_id):
    """Display Word document"""
    try:
        doc = Document(file)
        
        # Document stats
        col1, col2, col3 = st.columns(3)
        with col1:
            paragraphs = len([p for p in doc.paragraphs if p.text.strip()])
            st.metric("Paragraphs", paragraphs)
        with col2:
            st.metric("Tables", len(doc.tables))
        with col3:
            word_count = sum(len(p.text.split()) for p in doc.paragraphs)
            st.metric("Words", word_count)
        
        # Search
        search_term = st.text_input("🔍 Search in document", "", key=f"docx_search_{tab_id}")
        
        # Display content
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text = para.text
                if search_term:
                    text = re.sub(f'({re.escape(search_term)})', r'**\1**', text, flags=re.IGNORECASE)
                full_text.append(text)
        
        st.text_area("Document Content", '\n\n'.join(full_text), height=400, key=f"docx_content_{tab_id}")
        
        # Display tables
        if doc.tables:
            with st.expander("📊 Tables in Document"):
                for i, table in enumerate(doc.tables):
                    st.markdown(f"**Table {i+1}**")
                    data = []
                    for row in table.rows:
                        data.append([cell.text for cell in row.cells])
                    if data:
                        df = pd.DataFrame(data)
                        st.dataframe(df, use_container_width=True)
    
    except Exception as e:
        st.error(f"Error reading Word document: {str(e)}")

def display_pptx(file, tab_id):
    """Display PowerPoint presentation"""
    try:
        prs = Presentation(file)
        
        st.metric("Total Slides", len(prs.slides))
        
        # Slide navigation
        slide_num = st.number_input(
            "Select Slide",
            min_value=1,
            max_value=len(prs.slides),
            value=1,
            key=f"ppt_slide_{tab_id}"
        )
        
        slide = prs.slides[slide_num - 1]
        
        # Extract slide content
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        if slide_text:
            st.markdown(f"**Slide {slide_num} Content:**")
            for text in slide_text:
                st.info(text[:300] + "..." if len(text) > 300 else text)
        else:
            st.info("No text content on this slide")
        
        # Slide notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                with st.expander("📝 Slide Notes"):
                    st.write(notes)
    
    except Exception as e:
        st.error(f"Error reading PowerPoint: {str(e)}")

def display_excel(file, tab_id):
    """Display Excel file"""
    try:
        excel_file = pd.ExcelFile(file)
        
        # Sheet selection
        sheet_names = excel_file.sheet_names
        selected_sheet = st.selectbox(
            "Select Sheet",
            sheet_names,
            key=f"excel_sheet_{tab_id}"
        )
        
        # Read the selected sheet
        df = pd.read_excel(file, sheet_name=selected_sheet)
        
        # Data stats
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Rows", df.shape[0])
        with col2:
            st.metric("Columns", df.shape[1])
        with col3:
            st.metric("Memory", f"{df.memory_usage(deep=True).sum() / 1024:.1f} KB")
        
        # Display data
        st.dataframe(df, use_container_width=True, height=400)
        
        # Data summary
        with st.expander("📊 Data Summary"):
            st.write(df.describe())
        
        # Column info
        with st.expander("📋 Column Information"):
            for col in df.columns:
                st.text(f"{col}: {df[col].dtype} - {df[col].count()} values")
    
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")

def display_csv(file, tab_id):
    """Display CSV file"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        df = None
        error_msg = ""
        
        for encoding in encodings:
            try:
                file.seek(0)
                df = pd.read_csv(file, encoding=encoding, nrows=1000)  # Preview first 1000 rows
                break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                error_msg = str(e)
                continue
        
        if df is not None:
            # Data stats
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Rows", df.shape[0])
            with col2:
                st.metric("Columns", df.shape[1])
            
            # Display data
            st.dataframe(df, use_container_width=True, height=400)
            
            # Column info
            with st.expander("📋 Column Information"):
                for col in df.columns:
                    st.text(f"{col}: {df[col].dtype}")
        else:
            st.error(f"Could not read CSV file: {error_msg}")
    
    except Exception as e:
        st.error(f"Error reading CSV file: {str(e)}")

def display_text(file, tab_id):
    """Display text file"""
    try:
        content = file.getvalue().decode("utf-8")
        
        # Text stats
        lines = content.split('\n')
        words = content.split()
        chars = len(content)
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Lines", len(lines))
        with col2:
            st.metric("Words", len(words))
        with col3:
            st.metric("Characters", chars)
        
        # Search
        search_term = st.text_input("🔍 Search in text", "", key=f"txt_search_{tab_id}")
        
        # Display with line numbers
        show_line_numbers = st.checkbox("Show Line Numbers", True, key=f"txt_ln_{tab_id}")
        
        if show_line_numbers:
            numbered_content = ""
            for i, line in enumerate(lines, 1):
                line_text = line
                if search_term:
                    line_text = re.sub(f'({re.escape(search_term)})', r'**\1**', line_text, flags=re.IGNORECASE)
                numbered_content += f"{i:4d} | {line_text}\n"
            st.text_area("File Content", numbered_content, height=400, key=f"txt_content_{tab_id}")
        else:
            text_content = content
            if search_term:
                text_content = re.sub(f'({re.escape(search_term)})', r'**\1**', text_content, flags=re.IGNORECASE)
            st.text_area("File Content", text_content, height=400, key=f"txt_content_{tab_id}")
    
    except UnicodeDecodeError:
        st.error("Cannot display this file as text. Try downloading it instead.")
    except Exception as e:
        st.error(f"Error reading text file: {str(e)}")

def display_image(file, tab_id):
    """Display image file"""
    try:
        from PIL import Image, ExifTags
        
        img = Image.open(file)
        
        # Image info
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Format", img.format or "Unknown")
        with col2:
            st.metric("Size", f"{img.size[0]}x{img.size[1]}")
        with col3:
            st.metric("Mode", img.mode)
        
        # Display image
        st.image(img, caption=file.name, use_container_width=True)
        
        # EXIF data
        try:
            exif = img._getexif()
            if exif:
                with st.expander("📷 Image Metadata (EXIF)"):
                    for tag_id, value in exif.items():
                        tag = ExifTags.TAGS.get(tag_id, tag_id)
                        if value and str(value).strip():
                            st.text(f"{tag}: {value}")
        except:
            st.caption("No EXIF data available")
    
    except Exception as e:
        st.error(f"Error displaying image: {str(e)}")